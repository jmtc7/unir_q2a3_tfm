# -*- coding: utf-8 -*-
"""
Generador de la presentación de defensa del TFM "OpsInsight Analytics".
Máster en Análisis y Visualización de Datos Masivos (UNIR) · RedProyectum · Sopra Steria.
Autores: Efrel A. López Cáceres · José Miguel Torres Cámara · Víctor Romero Pardeiro.

Diseño estandarizado (ver GUIA_ESTILO.md): paleta industrial + teal + tipografía Segoe UI.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG = os.path.join(ROOT, "imgs")
OUT = os.path.join(ROOT, "presentacion", "OpsInsight_Analytics_Defensa_TFM.pptx")

# ----------------------------------------------------------------------------
# PALETA
# ----------------------------------------------------------------------------
INK      = RGBColor(0x0E, 0x1B, 0x2A)   # navy industrial (fondos oscuros)
PANEL    = RGBColor(0x12, 0x25, 0x3A)   # panel oscuro
TEAL     = RGBColor(0x16, 0xB8, 0xA6)   # acento principal
TEAL_D   = RGBColor(0x0E, 0x8C, 0x7F)   # teal oscuro
AMBER    = RGBColor(0xF5, 0xA6, 0x23)   # alertas / criticidad
RED      = RGBColor(0xE5, 0x48, 0x4D)   # errores
GREEN    = RGBColor(0x2F, 0xB5, 0x74)   # OK / completo
BG       = RGBColor(0xF4, 0xF7, 0xFA)   # fondo claro
CARD     = RGBColor(0xFF, 0xFF, 0xFF)   # tarjeta
INKTXT   = RGBColor(0x1A, 0x2A, 0x3A)   # texto principal
MUTED    = RGBColor(0x5B, 0x6B, 0x7B)   # texto secundario
HAIR     = RGBColor(0xDC, 0xE4, 0xEC)   # líneas finas
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TEALSOFT = RGBColor(0xE4, 0xF6, 0xF3)   # teal muy claro (fondos de chip)
INKSOFT  = RGBColor(0x1B, 0x30, 0x47)

# TIPOGRAFÍA
F_HEAD = "Segoe UI Semibold"
F_LIGHT = "Segoe UI Light"
F_BODY = "Segoe UI"
F_MONO = "Consolas"

# GEOMETRÍA (pulgadas)
SW, SH = 13.333, 7.5
ML, MR = 0.7, 0.7
CW = SW - ML - MR

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# HELPERS
# ----------------------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0, wrap=True):
    """runs: lista de párrafos; cada párrafo = lista de (texto, size, color, bold, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, fn) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = fn
    return tb


def R(t, sz, col, bold=False, fn=F_BODY):
    return (t, sz, col, bold, fn)


def pic_fit(s, path, x, y, w, h, ah="center", av="middle"):
    iw, ih = Image.open(path).size
    r = iw / ih; br = w / h
    if r > br:
        nw = w; nh = w / r
    else:
        nh = h; nw = h * r
    nx = x + (w - nw) / 2 if ah == "center" else (x if ah == "left" else x + (w - nw))
    ny = y + (h - nh) / 2 if av == "middle" else (y if av == "top" else y + (h - nh))
    return s.shapes.add_picture(path, Inches(nx), Inches(ny), Inches(nw), Inches(nh))


def footer(s, section, page):
    rect(s, ML, 7.02, CW, 0.012, fill=HAIR)
    txt(s, ML, 7.08, 7.5, 0.3,
        [[R("OpsInsight Analytics", 8.5, MUTED, True, F_HEAD),
          R("   ·   Defensa TFM · UNIR", 8.5, MUTED, False, F_BODY)]],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, SW - MR - 4.5, 7.08, 4.5, 0.3,
        [[R(section + "   ", 8.5, TEAL_D, True, F_HEAD),
          R("· %02d" % page, 8.5, MUTED, False, F_BODY)]],
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def content(title, kicker, section, page):
    """Slide de contenido estándar. Devuelve el slide."""
    s = slide()
    rect(s, 0, 0, SW, SH, fill=BG)
    # kicker
    rect(s, ML, 0.62, 0.28, 0.28, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.28)
    txt(s, ML + 0.42, 0.60, CW, 0.32,
        [[R(kicker.upper(), 11, TEAL_D, True, F_HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    # título
    txt(s, ML, 0.98, CW, 0.72, [[R(title, 27, INKTXT, True, F_HEAD)]],
        anchor=MSO_ANCHOR.MIDDLE)
    # regla de acento
    rect(s, ML, 1.72, 1.7, 0.05, fill=TEAL)
    footer(s, section, page)
    return s


def divider(number, title, subtitle, page):
    s = slide()
    rect(s, 0, 0, SW, SH, fill=INK)
    rect(s, 0, 0, 0.22, SH, fill=TEAL)
    # número gigante
    txt(s, ML + 0.1, 1.5, 6.5, 3.2, [[R(number, 200, INKSOFT, True, F_HEAD)]],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, ML + 0.35, 3.55, 10.5, 0.4, [[R("BLOQUE " + number, 13, TEAL, True, F_HEAD)]])
    txt(s, ML + 0.3, 3.95, 11.0, 1.1, [[R(title, 40, WHITE, True, F_HEAD)]])
    txt(s, ML + 0.32, 5.05, 10.5, 0.6, [[R(subtitle, 14, RGBColor(0x9F, 0xB2, 0xC2), False, F_BODY)]])
    txt(s, SW - MR - 3.0, 7.02, 3.0, 0.3,
        [[R("OpsInsight Analytics · %02d" % page, 9, RGBColor(0x6C, 0x80, 0x92), False, F_BODY)]],
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    return s


def chip(s, x, y, w, h, label, sub=None, fill=CARD, edge=TEAL, tcol=INKTXT):
    rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    rect(s, x, y, 0.09, h, fill=edge, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    if sub:
        txt(s, x + 0.24, y + 0.12, w - 0.35, h - 0.2,
            [[R(label, 12.5, tcol, True, F_HEAD)],
             [R(sub, 9.5, MUTED, False, F_BODY)]], anchor=MSO_ANCHOR.MIDDLE, space_after=2)
    else:
        txt(s, x + 0.24, y, w - 0.35, h,
            [[R(label, 12.5, tcol, True, F_HEAD)]], anchor=MSO_ANCHOR.MIDDLE)


def card(s, x, y, w, h, fill=CARD, edge=None):
    sp = rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    if edge:
        rect(s, x, y, w, 0.07, fill=edge, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    return sp


def bullets(s, x, y, w, h, items, size=13, gap=7, lead=1.12, bcol=TEAL, tcol=INKTXT):
    """items: lista de (texto, nivel). nivel 0/1."""
    paras = []
    for it in items:
        if isinstance(it, tuple):
            t, lvl = it
        else:
            t, lvl = it, 0
        if lvl == 0:
            paras.append([R("▸  ", size, bcol, True, F_HEAD), R(t, size, tcol, False, F_BODY)])
        else:
            paras.append([R("      –  ", size - 1, MUTED, False, F_BODY),
                          R(t, size - 1, MUTED, False, F_BODY)])
    txt(s, x, y, w, h, paras, space_after=gap, line_spacing=lead)


# ============================================================================
# SLIDE 1 — PORTADA
# ============================================================================
s = slide()
rect(s, 0, 0, SW, SH, fill=INK)
rect(s, 0, 0, SW, 0.22, fill=TEAL)
rect(s, 0, SH - 0.22, SW, 0.22, fill=TEAL)
# marca superior
txt(s, ML, 0.7, CW, 0.4,
    [[R("UNIVERSIDAD INTERNACIONAL DE LA RIOJA", 12, TEAL, True, F_HEAD),
      R("   ·   Escuela Superior de Ingeniería y Tecnología", 12, RGBColor(0xAF, 0xC0, 0xCE), False, F_BODY)]])
txt(s, ML, 1.12, CW, 0.35,
    [[R("Máster en Análisis y Visualización de Datos Masivos", 12.5, RGBColor(0xC6, 0xD3, 0xDD), False, F_BODY)]])
# título producto
txt(s, ML, 2.15, CW, 1.2, [[R("OpsInsight Analytics", 54, WHITE, True, F_HEAD)]])
rect(s, ML + 0.03, 3.35, 2.4, 0.06, fill=TEAL)
# subtítulo (título oficial)
txt(s, ML, 3.6, 11.4, 1.1,
    [[R("Diseño y desarrollo de una plataforma de monitorización operativa mediante "
        "cuadros de mando conectados a sistemas OCR y automatización con agentes",
        18, RGBColor(0xD7, 0xE1, 0xE9), False, F_LIGHT)]], line_spacing=1.12)
# autores
txt(s, ML, 5.05, 11.4, 0.4,
    [[R("Efrel Armando López Cáceres     ·     José Miguel Torres Cámara     ·     Víctor Romero Pardeiro",
        13.5, WHITE, True, F_HEAD)]])
txt(s, ML, 5.5, 11.4, 0.35,
    [[R("Director: ", 11.5, TEAL, True, F_HEAD), R("Rafael Martínez Ranera", 11.5, RGBColor(0xC6, 0xD3, 0xDD), False, F_BODY),
      R("        Tipo de trabajo: ", 11.5, TEAL, True, F_HEAD), R("Desarrollo", 11.5, RGBColor(0xC6, 0xD3, 0xDD), False, F_BODY)]])
# badge RedProyectum / Sopra Steria
rect(s, ML, 6.15, 4.15, 0.5, fill=PANEL, line=TEAL_D, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.4)
txt(s, ML + 0.22, 6.15, 3.9, 0.5,
    [[R("RedProyectum", 11, TEAL, True, F_HEAD), R("  ·  en colaboración con Sopra Steria", 10.5, RGBColor(0xC6, 0xD3, 0xDD), False, F_BODY)]],
    anchor=MSO_ANCHOR.MIDDLE)
txt(s, SW - MR - 3.5, 6.15, 3.5, 0.5,
    [[R("Defensa de Trabajo Fin de Máster", 10.5, RGBColor(0x9F, 0xB2, 0xC2), False, F_BODY)]],
    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "PRESENTA: Efrel (apertura del equipo). ~40 s.\n"
         "Buenos días. Somos Efrel, José Miguel y Víctor. Presentamos OpsInsight Analytics, "
         "una plataforma local multiagente que transforma documentación operativa heterogénea "
         "en inteligencia analítica accionable. TFM en modalidad RedProyectum con Sopra Steria, "
         "dirigido por Rafael Martínez Ranera.")

# ============================================================================
# SLIDE 2 — AGENDA
# ============================================================================
s = content("Agenda de la defensa", "Índice", "Agenda", 2)
items = [
    ("01", "Contexto y objetivos", "El problema, el objetivo general y los 9 objetivos específicos; metodología CRISP-DM.", TEAL),
    ("02", "Arquitectura de la solución", "Las 5 capas, el stack 100% local y el workflow integrado en n8n.", AMBER),
    ("03", "Módulos del pipeline", "Ingesta + OCR, validación en 2 capas, agente RAG, registro y dashboard.", TEAL),
    ("04", "Resultados y conclusiones", "Grado de cumplimiento, limitaciones, trabajo futuro y conclusiones.", GREEN),
]
y = 2.15
for num, tt, ds, col in items:
    card(s, ML, y, CW, 1.02, fill=CARD)
    rect(s, ML, y, 0.12, 1.02, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    txt(s, ML + 0.35, y, 1.3, 1.02, [[R(num, 34, col, True, F_HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, ML + 1.7, y + 0.14, CW - 2.0, 0.8,
        [[R(tt, 16, INKTXT, True, F_HEAD)], [R(ds, 11, MUTED, False, F_BODY)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=3)
    y += 1.16
notes(s, "PRESENTA: Efrel. ~30 s.\nRecorremos cuatro bloques: contexto y objetivos, arquitectura, "
         "los módulos del pipeline y, para cerrar, resultados y conclusiones.")

# ============================================================================
# SLIDE 3 — DIVIDER 01
# ============================================================================
s = divider("01", "Contexto y objetivos", "Por qué el conocimiento operativo está fragmentado y qué nos propusimos resolver.", 3)
notes(s, "PRESENTA: Efrel.")

# ============================================================================
# SLIDE 4 — CONTEXTO Y PROBLEMA
# ============================================================================
s = content("El conocimiento operativo está fragmentado", "Contexto y problema", "Contexto", 4)
# tarjeta grande de problema (izquierda)
card(s, ML, 2.15, 6.6, 4.4, fill=CARD)
txt(s, ML + 0.35, 2.45, 6.0, 0.5, [[R("El problema de partida", 15, INKTXT, True, F_HEAD)]])
bullets(s, ML + 0.35, 3.05, 6.0, 3.3, [
    "Fuentes heterogéneas: formularios online, partes en PDF, reportes manuscritos y manuales técnicos.",
    "Nomenclaturas inconsistentes, formatos dispares y campos incompletos.",
    "Sin una visión consolidada del comportamiento de los activos.",
    "Consecuencia: la toma de decisiones de mantenimiento se hace sin evidencia.",
], size=12.5, gap=9)
# columna derecha: contexto + meta
card(s, ML + 6.9, 2.15, CW - 6.9, 2.05, fill=INK)
txt(s, ML + 7.15, 2.4, CW - 7.35, 1.7,
    [[R("Entorno de referencia", 11, TEAL, True, F_HEAD)],
     [R("Mantenimiento industrial en organizaciones de servicios tecnológicos como Sopra Steria.",
        12, WHITE, False, F_BODY)],
     [R("Datos: AI4I 2020 y MetroPT-3 (compresor de aire de un metro).", 11, RGBColor(0xC6, 0xD3, 0xDD), False, F_BODY)]],
    space_after=6, line_spacing=1.1)
card(s, ML + 6.9, 4.5, CW - 6.9, 2.05, fill=TEALSOFT, edge=TEAL)
txt(s, ML + 7.15, 4.75, CW - 7.35, 1.7,
    [[R("Nuestra meta", 11, TEAL_D, True, F_HEAD)],
     [R("Pasar de documentación heterogénea a inteligencia analítica accionable:", 12, INKTXT, True, F_HEAD)],
     [R("local, trazable y reproducible.", 12, TEAL_D, True, F_HEAD)]],
    space_after=6, line_spacing=1.1)
notes(s, "PRESENTA: Efrel. ~1 min.\nEl conocimiento del negocio queda disperso en papeles, PDFs y sistemas "
         "distintos. No hay dependencia de datos reales de empresa: usamos datasets públicos de referencia "
         "(AI4I 2020 y MetroPT-3) sugeridos por Sopra Steria. El objetivo es cerrar la cadena del dato: de "
         "la ingesta a la decisión.")

# ============================================================================
# SLIDE 5 — OBJETIVOS
# ============================================================================
s = content("Objetivo general y objetivos específicos", "Objetivos", "Objetivos", 5)
# objetivo general
card(s, ML, 2.05, CW, 1.05, fill=INK)
rect(s, ML, 2.05, 0.12, 1.05, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
txt(s, ML + 0.35, 2.05, CW - 0.6, 1.05,
    [[R("OBJETIVO GENERAL   ", 11, TEAL, True, F_HEAD),
      R("Diseñar, implementar y validar OpsInsight Analytics, una plataforma multiagente local que "
        "transforme información operativa dispersa y heterogénea en analítica estructurada, trazable y accionable.",
        12.5, WHITE, False, F_BODY)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
oe = [
    ("OE1", "Analizar el dominio y los datasets de referencia (AI4I 2020, MetroPT-3)."),
    ("OE2", "Comparar soluciones OCR (Tesseract vs Paddle)."),
    ("OE3", "Módulo de ingesta multiformato → JSON estandarizado."),
    ("OE4", "Validación del JSON en dos capas (estructural + semántica)."),
    ("OE5", "Repositorio analítico consolidado en MySQL."),
    ("OE6", "Análisis descriptivo y temporal del repositorio."),
    ("OE7", "Índice de criticidad (frecuencia + parada + coste)."),
    ("OE8", "Agente de diagnóstico técnico mediante RAG."),
    ("OE9", "Cuadro de mando interactivo en Grafana."),
]
colw = (CW - 0.4) / 2
x0 = ML; x1 = ML + colw + 0.4
y = 3.35
for i, (code, desc) in enumerate(oe):
    col = x0 if i < 5 else x1
    row = i if i < 5 else i - 5
    yy = y + row * 0.66
    txt(s, col, yy, colw, 0.6,
        [[R(code + "  ", 12, TEAL_D, True, F_HEAD), R(desc, 11, INKTXT, False, F_BODY)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    rect(s, col, yy + 0.6, colw, 0.010, fill=HAIR)
notes(s, "PRESENTA: Efrel. ~1 min.\nNueve objetivos específicos, agrupados en tres bloques: preparación del "
         "dato (OE1-OE5), analítica (OE6-OE7) y las dos capas inteligentes (OE8 RAG, OE9 dashboard). "
         "Volveremos sobre el grado de cumplimiento al final.")

# ============================================================================
# SLIDE 6 — METODOLOGÍA CRISP-DM
# ============================================================================
s = content("Metodología: CRISP-DM adaptado", "Método y equipo", "Metodología", 6)
fases = [
    ("1", "Comprensión\ndel negocio", "Alcance, indicadores\ny fuentes fiables", TEAL),
    ("2", "Comprensión\nde los datos", "Manual sintético +\nplantilla · OE1", TEAL),
    ("3", "Preparación\nde los datos", "Pipeline OCR +\nagente extractor", AMBER),
    ("4", "Modelado", "Modelo MySQL +\ncriticidad · OE6-7", AMBER),
    ("5", "Evaluación", "OCR, esquema y\nusabilidad", TEAL),
    ("6", "Despliegue", "Integración +\nDocker + TFM", GREEN),
]
n = len(fases); gap = 0.2
cwid = (CW - gap * (n - 1)) / n
x = ML; y = 2.25
for num, tt, ds, col in fases:
    card(s, x, y, cwid, 2.2, fill=CARD)
    rect(s, x, y, cwid, 0.5, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    rect(s, x, y + 0.25, cwid, 0.25, fill=col)
    txt(s, x, y + 0.02, cwid, 0.5, [[R("FASE " + num, 10.5, WHITE, True, F_HEAD)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 0.1, y + 0.62, cwid - 0.2, 0.9,
        [[R(l, 11.5, INKTXT, True, F_HEAD)] for l in tt.split("\n")],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
    txt(s, x + 0.1, y + 1.5, cwid - 0.2, 0.65,
        [[R(l, 9, MUTED, False, F_BODY)] for l in ds.split("\n")],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, space_after=0, line_spacing=1.0)
    x += cwid + gap
# reparto de responsabilidades
card(s, ML, 4.75, CW, 1.75, fill=BG, edge=TEAL)
rect(s, ML, 4.75, CW, 0.07, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
txt(s, ML + 0.3, 4.9, CW - 0.6, 0.4, [[R("Reparto de responsabilidades (16 semanas)", 13, INKTXT, True, F_HEAD)]])
repw = (CW - 0.6 - 0.4 * 2) / 3
rx = ML + 0.3
for name, tasks in [
    ("Efrel A. López", "Modelo de datos MySQL · análisis temporal y criticidad (OE6-7) · dashboard (OE9) · integración"),
    ("José Miguel Torres", "Pipeline OCR (Tesseract/Paddle) · ingesta y extracción · validación · registro (OE2-5)"),
    ("Víctor Romero", "Análisis descriptivo · índice de criticidad · agente de diagnóstico RAG (OE6-8)"),
]:
    txt(s, rx, 5.35, repw, 1.05,
        [[R(name, 11.5, TEAL_D, True, F_HEAD)], [R(tasks, 9.5, MUTED, False, F_BODY)]],
        space_after=3, line_spacing=1.05)
    rx += repw + 0.4
notes(s, "PRESENTA: Efrel. ~50 s.\nSeguimos CRISP-DM, adaptado al dominio industrial, en seis fases sobre "
         "16 semanas. El trabajo se repartió por afinidad técnica pero con revisión cruzada: todos conocemos "
         "todo el proyecto.")

# ============================================================================
# SLIDE 7 — DIVIDER 02
# ============================================================================
s = divider("02", "Arquitectura de la solución", "Cinco capas técnicas, stack 100% local y orquestación con n8n.", 7)
notes(s, "PRESENTA: Víctor.")

# ============================================================================
# SLIDE 8 — ARQUITECTURA GENERAL
# ============================================================================
s = content("Arquitectura: cinco capas encadenadas", "Visión general", "Arquitectura", 8)
capas = [
    ("Ingesta", "Fuentes\nheterogéneas", TEAL),
    ("Extracción", "OCR + LLM\n→ JSON", TEAL),
    ("Validación", "Estructural +\nsemántica", AMBER),
    ("Registro", "MySQL +\nCSV", TEAL_D),
    ("Visualización", "Dashboard\nGrafana", GREEN),
]
n = len(capas); gap = 0.72
cwid = (CW - gap * (n - 1)) / n
x = ML; y = 2.15
for i, (tt, ds, col) in enumerate(capas):
    card(s, x, y, cwid, 1.55, fill=CARD)
    rect(s, x, y, cwid, 0.09, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    txt(s, x + 0.1, y + 0.2, cwid - 0.2, 0.5, [[R(tt, 13, INKTXT, True, F_HEAD)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 0.1, y + 0.72, cwid - 0.2, 0.7,
        [[R(l, 10, MUTED, False, F_BODY)] for l in ds.split("\n")],
        align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.0)
    if i < n - 1:
        txt(s, x + cwid - 0.02, y, gap, 1.55, [[R("→", 22, TEAL, True, F_HEAD)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += cwid + gap
# diagrama concreto NATIVO (coherente: termina en Grafana, no Streamlit)
y0 = 4.15
bh = 1.5


def _node(x, w, title, sub, edge, fill=CARD, tcol=INKTXT):
    card(s, x, y0, w, bh, fill=fill)
    rect(s, x, y0, w, 0.08, fill=edge, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    paras = [[R(title, 12.5, tcol, True, F_HEAD)]] + \
            [[R(l, 9.5, MUTED, False, F_BODY)] for l in sub.split("\n")]
    txt(s, x + 0.1, y0 + 0.18, w - 0.2, bh - 0.3, paras,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=3, line_spacing=1.0)


def _arrow(x, w):
    txt(s, x, y0, w, bh, [[R("→", 20, TEAL, True, F_HEAD)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


gp = 0.3
xe = ML
_node(xe, 1.95, "Entradas", "Formulario · PDF\nImagen (OCR)", TEAL); xe += 1.95
_arrow(xe, gp); xe += gp
# contenedor n8n con dos agentes
nw = 4.95
card(s, xe, y0, nw, bh, fill=TEALSOFT)
rect(s, xe, y0, nw, 0.08, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
txt(s, xe + 0.1, y0 + 0.13, nw - 0.2, 0.3,
    [[R("n8n · orquestación de agentes", 10.5, TEAL_D, True, F_HEAD)]], align=PP_ALIGN.CENTER)
pad = 0.18
chw = (nw - pad * 2 - 0.24) / 2
cy = y0 + 0.5
cx = xe + pad
rect(s, cx, cy, chw, 0.82, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
txt(s, cx + 0.05, cy + 0.06, chw - 0.1, 0.7,
    [[R("Agente de ingesta", 10.5, INKTXT, True, F_HEAD)], [R("→ JSON estandarizado", 8.5, MUTED, False, F_BODY)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=1)
txt(s, cx + chw, cy, 0.24, 0.82, [[R("→", 13, TEAL_D, True, F_HEAD)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, cx + chw + 0.24, cy, chw, 0.82, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
txt(s, cx + chw + 0.29, cy + 0.06, chw - 0.1, 0.7,
    [[R("Agente de análisis", 10.5, INKTXT, True, F_HEAD)], [R("RAG + manual MetroPT-3", 8.5, MUTED, False, F_BODY)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=1)
xe += nw
_arrow(xe, gp); xe += gp
_node(xe, 1.95, "Registro", "MySQL + CSV", TEAL_D); xe += 1.95
_arrow(xe, gp); xe += gp
_node(xe, 1.85, "Grafana", "Cuadro de mando", GREEN)
txt(s, ML, 5.85, CW, 0.3,
    [[R("Orquestación jerárquica de agentes coordinados por n8n · flujo trazable de extremo a extremo",
        11, MUTED, False, F_BODY)]], align=PP_ALIGN.CENTER)
notes(s, "PRESENTA: Víctor. ~1 min.\nEl dato recorre cinco capas encadenadas. Cada módulo es independiente y "
         "sustituible sin rediseñar el conjunto — esa modularidad es clave para la trazabilidad. Todo lo "
         "orquesta n8n.")

# ============================================================================
# SLIDE 9 — STACK TECNOLÓGICO
# ============================================================================
s = content("Stack tecnológico: 100% local", "Tecnologías", "Arquitectura", 9)
stack = [
    ("n8n", "Orquestación de agentes y workflow", TEAL),
    ("Ollama", "LLM local: Qwen2.5:3b · llama3.2:3b", TEAL),
    ("Tesseract · PaddleOCR", "Reconocimiento óptico de caracteres", AMBER),
    ("RAG · nomic-embed-text", "Recuperación semántica del manual", AMBER),
    ("MySQL 8.0", "Repositorio analítico de incidencias", TEAL_D),
    ("Grafana 11.4", "Cuadro de mando auto-provisionado", GREEN),
    ("Docker Compose", "Despliegue reproducible del stack", MUTED),
    ("Python · node:test", "Servicio OCR (FastAPI) y validador testeado", MUTED),
]
colw = (CW - 0.4) / 2
x0 = ML; x1 = ML + colw + 0.4
y = 2.2
for i, (name, ds, col) in enumerate(stack):
    col_x = x0 if i % 2 == 0 else x1
    row = i // 2
    yy = y + row * 0.92
    chip(s, col_x, yy, colw, 0.76, name, ds, fill=CARD, edge=col)
# franja privacidad
card(s, ML, 6.05, CW, 0.62, fill=INK)
txt(s, ML + 0.3, 6.05, CW - 0.6, 0.62,
    [[R("100% local  ", 13, TEAL, True, F_HEAD),
      R("Sin dependencias en la nube ni exposición de datos sensibles → privacidad y reproducibilidad garantizadas.",
        12, WHITE, False, F_BODY)]], anchor=MSO_ANCHOR.MIDDLE)
notes(s, "PRESENTA: Víctor. ~50 s.\nTodo el stack corre en local con Docker. El punto diferencial: LLMs "
         "de 3B parámetros ejecutables en CPU con pocos GB de RAM. Nada sale a la nube — privacidad total, "
         "algo crítico para un entorno como Sopra Steria.")

# ============================================================================
# SLIDE 10 — WORKFLOW N8N
# ============================================================================
s = content("Workflow integrado en n8n", "Orquestación", "Arquitectura", 10)
card(s, ML, 2.1, CW, 3.35, fill=WHITE)
pic_fit(s, os.path.join(IMG, "n8n", "0_global.png"), ML + 0.25, 2.25, CW - 0.5, 3.05)
zonas = [
    ("Ingesta documental", "Formulario · PDF · OCR", TEAL),
    ("Validación del JSON", "Estructural + semántica", AMBER),
    ("Análisis documental", "RAG + AI Agent", TEAL_D),
    ("Registro", "MySQL + CSV", GREEN),
]
n = len(zonas); gap = 0.25
cwid = (CW - gap * (n - 1)) / n
x = ML; y = 5.65
for tt, ds, col in zonas:
    rect(s, x, y, 0.1, 0.85, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    txt(s, x + 0.22, y, cwid - 0.22, 0.85,
        [[R(tt, 11.5, INKTXT, True, F_HEAD)], [R(ds, 9.5, MUTED, False, F_BODY)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=2)
    x += cwid + gap
notes(s, "PRESENTA: Víctor. ~40 s.\nEste es el workflow completo en n8n, organizado en cuatro zonas de "
         "responsabilidad. Ahora entramos en cada módulo en detalle.")

# ============================================================================
# SLIDE 11 — DIVIDER 03
# ============================================================================
s = divider("03", "Módulos del pipeline", "Del reporte crudo al diagnóstico registrado: ingesta, validación, RAG, registro y dashboard.", 11)
notes(s, "PRESENTA: José Miguel.")

# ============================================================================
# SLIDE 12 — INGESTA + OCR
# ============================================================================
s = content("Ingesta omnicanal y extracción OCR/LLM", "Módulo · OE2-OE3", "Ingesta", 12)
vias = [
    ("Formulario online", "Entrada web estructurada", TEAL),
    ("PDF digital", "Extracción de texto + LLM", TEAL_D),
    ("Foto (impreso / manuscrito)", "OCR → texto → LLM", AMBER),
]
y = 2.1
for tt, ds, col in vias:
    chip(s, ML, y, 6.3, 0.8, tt, ds, fill=CARD, edge=col)
    y += 0.92
# salida
card(s, ML, 4.86, 6.3, 1.55, fill=INK)
txt(s, ML + 0.28, 5.05, 5.9, 1.3,
    [[R("Salida unificada", 11, TEAL, True, F_HEAD)],
     [R("JSON de incidencia estandarizado", 13, WHITE, True, F_HEAD)],
     [R("Extracción con Qwen2.5:3b (respeta formatos estrictos, no 'interpreta').", 10.5, RGBColor(0xC6, 0xD3, 0xDD), False, F_BODY)]],
    space_after=5, line_spacing=1.05)
# comparativa OCR (derecha)
card(s, ML + 6.6, 2.1, CW - 6.6, 4.31, fill=CARD)
txt(s, ML + 6.9, 2.32, CW - 7.1, 0.4, [[R("Comparativa OCR (OE2)", 13.5, INKTXT, True, F_HEAD)]])
rows = [
    ("Tesseract", "Rápido y ligero, pero peor con manuscritos.", RED),
    ("PaddleOCR", "Más pesado, pero mayor precisión.", GREEN),
]
yy = 2.9
for name, ds, col in rows:
    rect(s, ML + 6.9, yy, 0.1, 0.9, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    txt(s, ML + 7.15, yy, CW - 7.45, 0.9,
        [[R(name, 12.5, INKTXT, True, F_HEAD)], [R(ds, 10.5, MUTED, False, F_BODY)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=2)
    yy += 1.0
card(s, ML + 6.9, 4.95, CW - 7.2, 1.35, fill=TEALSOFT, edge=GREEN)
txt(s, ML + 7.15, 5.1, CW - 7.6, 1.1,
    [[R("Decisión", 10.5, GREEN, True, F_HEAD)],
     [R("PaddleOCR", 15, TEAL_D, True, F_HEAD)],
     [R("con evidencia experimental sobre los formularios del Anexo A.", 10, MUTED, False, F_BODY)]],
    space_after=3, line_spacing=1.02)
notes(s, "PRESENTA: José Miguel. ~1:15 min.\nTres vías de entrada convergen en un único JSON. La extracción "
         "usa Qwen2.5 porque respeta formatos estrictos sin inventar. Para el OCR comparamos Tesseract y "
         "Paddle sobre formularios reales del Anexo A: elegimos Paddle por precisión, sobre todo en "
         "manuscritos.")

# ============================================================================
# SLIDE 13 — VALIDACIÓN EN DOS CAPAS
# ============================================================================
s = content("Validación en dos capas: el gate de calidad", "Módulo · OE4", "Validación", 13)
# capa 1
card(s, ML, 2.1, 5.85, 1.35, fill=CARD, edge=RED)
txt(s, ML + 0.28, 2.28, 5.4, 1.1,
    [[R("Capa 1 · Estructural  ", 12.5, INKTXT, True, F_HEAD), R("(bloqueante)", 10.5, RED, True, F_HEAD)],
     [R("Campos requeridos, tipos y rangos físicamente posibles. Si falla → validacion_ok=0, no pasa al agente.",
        10.5, MUTED, False, F_BODY)]], space_after=4, line_spacing=1.05)
# capa 2
card(s, ML, 3.6, 5.85, 1.35, fill=CARD, edge=AMBER)
txt(s, ML + 0.28, 3.78, 5.4, 1.1,
    [[R("Capa 2 · Semántica  ", 12.5, INKTXT, True, F_HEAD), R("(informativa)", 10.5, AMBER, True, F_HEAD)],
     [R("Coherencia entre variables según el manual MetroPT-3. Inconsistencias → alertas_validacion (visibles en Grafana).",
        10.5, MUTED, False, F_BODY)]], space_after=4, line_spacing=1.05)
# calidad de software
card(s, ML, 5.1, 5.85, 1.3, fill=INK)
txt(s, ML + 0.28, 5.28, 5.4, 1.05,
    [[R("Código testeado = código en producción", 12, TEAL, True, F_HEAD)],
     [R("Función JS pura + 19 tests (node:test) → inyectada en el nodo Code de n8n vía build. Nunca divergen.",
        10.5, WHITE, False, F_BODY)]], space_after=4, line_spacing=1.05)
# tabla 5 reglas (derecha)
txt(s, ML + 6.1, 2.1, CW - 6.1, 0.35, [[R("5 reglas semánticas del manual", 13, INKTXT, True, F_HEAD)]])
reglas = [
    ("COMP_DV_ELECTRIC", "COMP y DV_Electric no simultáneas"),
    ("DV_BAJO_CARGA", "Bajo carga, presión DV ≈ 0 bar"),
    ("R_LEJOS_TP3", "Presión depósito R cercana a TP3"),
    ("LPS_INCOHERENTE", "LPS=1 ↔ R < 7 bar"),
    ("CORRIENTE_FUERA_RANGO", "Corriente ≈ nominales {0,4,7,9} A"),
]
yy = 2.55
rw = CW - 6.1
for i, (code, ds) in enumerate(reglas):
    fill = WHITE if i % 2 == 0 else BG
    card(s, ML + 6.1, yy, rw, 0.72, fill=fill)
    txt(s, ML + 6.32, yy + 0.08, rw - 0.4, 0.6,
        [[R(code, 10.5, TEAL_D, True, F_MONO)], [R(ds, 10, MUTED, False, F_BODY)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=1, line_spacing=1.0)
    yy += 0.78
notes(s, "PRESENTA: José Miguel. ~1:15 min.\nEntre ingesta y análisis metimos un gate de calidad en dos "
         "capas. La estructural es bloqueante; la semántica genera alertas de coherencia contra el manual. "
         "Detalle de ingeniería: la lógica es una función pura con 19 tests que se inyecta automáticamente "
         "en n8n, así el código testeado y el de producción nunca divergen.")

# ============================================================================
# SLIDE 14 — AGENTE RAG
# ============================================================================
s = content("Agente de diagnóstico documental (RAG)", "Módulo · OE8", "Análisis RAG", 14)
bullets(s, ML, 2.15, 6.2, 2.6, [
    "AI Agent (n8n) + llama3.2:3b vía Ollama.",
    "Tool buscador_manual: búsqueda vectorial sobre el manual MetroPT-3 indexado con embeddings nomic-embed-text.",
    "El JSON de incidencia se convierte en consulta en lenguaje natural.",
    "Grounded en el manual: recupera el fallo que mejor encaja, sin alucinar.",
], size=12.5, gap=9)
# salidas obligatorias
card(s, ML, 4.9, 6.2, 1.5, fill=INK)
txt(s, ML + 0.28, 5.05, 5.8, 0.35, [[R("Salida estructurada obligatoria", 11, TEAL, True, F_HEAD)]])
outs = [("Criticidad", "Alta / Media / Baja"), ("Tiempo est.", "horas"), ("Coste est.", "€")]
ox = ML + 0.28; ow = (6.2 - 0.56 - 0.3 * 2) / 3
for tt, ds in outs:
    rect(s, ox, 5.45, ow, 0.8, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    txt(s, ox, 5.5, ow, 0.7, [[R(tt, 11, WHITE, True, F_HEAD)], [R(ds, 9, TEAL, False, F_BODY)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=1)
    ox += ow + 0.3
# imagen análisis
card(s, ML + 6.5, 2.15, CW - 6.5, 4.25, fill=WHITE)
pic_fit(s, os.path.join(IMG, "n8n", "3_analisis.png"), ML + 6.7, 2.3, CW - 6.9, 3.6)
txt(s, ML + 6.5, 5.95, CW - 6.5, 0.35,
    [[R("Qwen (extracción estricta)  vs  llama (razonamiento y lenguaje natural)", 10, MUTED, False, F_BODY)]],
    align=PP_ALIGN.CENTER)
notes(s, "PRESENTA: Víctor. ~1:15 min.\nEl corazón inteligente: un AI Agent con un LLM local equipado con una "
         "herramienta de búsqueda vectorial sobre el manual. Traduce la incidencia a lenguaje natural, "
         "recupera el fallo que mejor encaja y devuelve criticidad, tiempo y coste — siempre anclado al "
         "manual, sin alucinar. Elegimos llama para razonar y Qwen para extraer.")

# ============================================================================
# SLIDE 15 — REGISTRO
# ============================================================================
s = content("Registro con visibilidad completa", "Módulo · OE5", "Registro", 15)
# dos destinos
card(s, ML, 2.15, 3.6, 1.5, fill=CARD, edge=TEAL)
txt(s, ML + 0.28, 2.35, 3.2, 1.2,
    [[R("MySQL", 15, INKTXT, True, F_HEAD)],
     [R("Tabla registro_incidencias: diagnóstico + campos de validación.", 10.5, MUTED, False, F_BODY)]],
    space_after=4, line_spacing=1.05)
card(s, ML + 3.85, 2.15, 3.05, 1.5, fill=CARD, edge=TEAL_D)
txt(s, ML + 4.13, 2.35, 2.6, 1.2,
    [[R("CSV", 15, INKTXT, True, F_HEAD)],
     [R("Copia en texto plano para consulta rápida.", 10.5, MUTED, False, F_BODY)]],
    space_after=4, line_spacing=1.05)
# todos los registros
txt(s, ML, 3.9, 6.9, 0.35, [[R("Todos los eventos quedan trazados:", 12.5, INKTXT, True, F_HEAD)]])
estados = [
    ("Válidos", "validacion_ok=1", GREEN),
    ("Rechazados", "criticidad = ERROR", AMBER),
    ("Fallos de sistema", "Error Handler → SISTEMA_ERROR", RED),
]
yy = 4.35
for tt, ds, col in estados:
    rect(s, ML, yy, 0.1, 0.6, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    txt(s, ML + 0.22, yy, 6.6, 0.6,
        [[R(tt + "  —  ", 12, INKTXT, True, F_HEAD), R(ds, 11, MUTED, False, F_MONO)]],
        anchor=MSO_ANCHOR.MIDDLE)
    yy += 0.66
# callout optimización
card(s, ML + 7.1, 2.15, CW - 7.1, 4.25, fill=INK)
rect(s, ML + 7.1, 2.15, CW - 7.1, 0.07, fill=AMBER, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
txt(s, ML + 7.4, 2.4, CW - 7.7, 0.7, [[R("Optimizaciones clave", 14, AMBER, True, F_HEAD)]])
bullets(s, ML + 7.4, 3.05, CW - 7.7, 3.2, [
    "Desacople del RAG: el manual se indexa una sola vez (antes se re-vectorizaba en cada incidencia → 30 min).",
    "Fin del \"efecto eco\": embudo lineal que elimina los registros duplicados del AI Agent.",
    "Ingesta modular: canales aislados evitan esquemas híbridos que confundían al agente.",
], size=11.5, gap=10, bcol=AMBER, tcol=WHITE)
notes(s, "PRESENTA: José Miguel. ~1 min.\nPersistimos en MySQL y CSV. Lo importante: TODO queda trazado, "
         "incluidos rechazos y fallos de infraestructura capturados por un error workflow. A la derecha, "
         "tres optimizaciones de rendimiento que hicimos sobre la versión base: la más grande, desacoplar "
         "el RAG bajó una incidencia de 30 minutos a segundos.")

# ============================================================================
# SLIDE 16 — DASHBOARD
# ============================================================================
s = content("Cuadro de mando en Grafana", "Módulo · OE9", "Dashboard", 16)
card(s, ML, 2.1, 8.2, 4.3, fill=WHITE)
pic_fit(s, os.path.join(IMG, "grafana", "dashboard.png"), ML + 0.15, 2.25, 7.9, 4.0)
# KPIs (derecha)
kpis = [
    ("Total incidencias", "COUNT(*)", TEAL),
    ("% Validación OK", "AVG(validacion_ok)", GREEN),
    ("Coste total (€)", "SUM(coste)", AMBER),
    ("MTTR (h)", "AVG(tiempo)", TEAL_D),
]
kx = ML + 8.5; kw = CW - 8.5
ky = 2.1
for tt, ds, col in kpis:
    card(s, kx, ky, kw, 0.72, fill=CARD)
    rect(s, kx, ky, 0.09, 0.72, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    txt(s, kx + 0.24, ky + 0.06, kw - 0.4, 0.6,
        [[R(tt, 12, INKTXT, True, F_HEAD)], [R(ds, 9.5, MUTED, False, F_MONO)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=1)
    ky += 0.82
card(s, kx, ky, kw, 1.0, fill=TEALSOFT, edge=TEAL)
txt(s, kx + 0.24, ky + 0.1, kw - 0.4, 0.8,
    [[R("Dashboards-as-code", 11.5, TEAL_D, True, F_HEAD)],
     [R("Auto-provisionado al arrancar · refresco 30 s · 9 paneles.", 10, MUTED, False, F_BODY)]],
    space_after=3, line_spacing=1.05)
notes(s, "PRESENTA: Efrel. ~1 min.\nAquí se hace visible el valor. El dashboard se auto-provisiona con el "
         "contenedor (dashboards-as-code), sin configuración manual, y refresca cada 30 s. Expone KPIs de "
         "negocio, estado de validación, alertas semánticas y el registro completo. Cierra la cadena del "
         "dato: de un reporte heterogéneo a una decisión informada.")

# ============================================================================
# SLIDE 17 — DIVIDER 04
# ============================================================================
s = divider("04", "Resultados y conclusiones", "Grado de cumplimiento, limitaciones honestas y hoja de ruta.", 17)
notes(s, "PRESENTA: Efrel.")

# ============================================================================
# SLIDE 18 — RESULTADOS / CUMPLIMIENTO
# ============================================================================
s = content("Resultados y grado de cumplimiento", "Evaluación", "Resultados", 18)
# KPIs cuantitativos
metrics = [("3", "vías de ingesta"), ("5", "reglas semánticas"), ("9", "paneles de KPI"), ("100%", "local")]
mw = (CW - 0.3 * 3) / 4
mx = ML
for val, lab in metrics:
    card(s, mx, 2.1, mw, 1.15, fill=INK)
    txt(s, mx, 2.2, mw, 0.7, [[R(val, 30, TEAL, True, F_HEAD)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, mx, 2.85, mw, 0.3, [[R(lab, 11, WHITE, False, F_BODY)]], align=PP_ALIGN.CENTER)
    mx += mw + 0.3
# validación end-to-end
txt(s, ML, 3.55, CW, 0.35,
    [[R("Validación de extremo a extremo", 13, INKTXT, True, F_HEAD),
      R("  sobre documentación sintética basada en AI4I 2020 y MetroPT-3.", 11.5, MUTED, False, F_BODY)]])
# matriz OE
txt(s, ML, 4.05, CW, 0.3, [[R("Cumplimiento de objetivos específicos", 12.5, INKTXT, True, F_HEAD)]])
oe_status = [("OE1", GREEN), ("OE2", GREEN), ("OE3", GREEN), ("OE4", GREEN), ("OE5", GREEN),
             ("OE6", AMBER), ("OE7", AMBER), ("OE8", GREEN), ("OE9", GREEN)]
ow = (CW - 0.2 * 8) / 9
ox = ML
for code, col in oe_status:
    rect(s, ox, 4.45, ow, 0.7, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    txt(s, ox, 4.45, ow, 0.7, [[R(code, 13, WHITE, True, F_HEAD)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ox += ow + 0.2
# leyenda
rect(s, ML, 5.4, 0.3, 0.3, fill=GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2)
txt(s, ML + 0.4, 5.4, 4.0, 0.3, [[R("Completo y validado (7/9)", 11, INKTXT, False, F_BODY)]], anchor=MSO_ANCHOR.MIDDLE)
rect(s, ML + 4.3, 5.4, 0.3, 0.3, fill=AMBER, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2)
txt(s, ML + 4.7, 5.4, 6.0, 0.3,
    [[R("Parcial — soportado descriptivamente (OE6 temporal, OE7 índice compuesto)", 11, INKTXT, False, F_BODY)]],
    anchor=MSO_ANCHOR.MIDDLE)
card(s, ML, 5.95, CW, 0.55, fill=TEALSOFT, edge=TEAL)
txt(s, ML + 0.3, 5.95, CW - 0.6, 0.55,
    [[R("Prototipo de producto mínimo viable, honesto en su alcance: de documentación heterogénea a "
        "inteligencia accionable, local y trazable.", 11.5, TEAL_D, True, F_HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
notes(s, "PRESENTA: Efrel. ~1 min.\nValidamos el sistema completo de extremo a extremo con datos sintéticos "
         "sobre los datasets de referencia. De los nueve objetivos, siete quedan completos y validados; OE6 "
         "y OE7 (análisis temporal e índice de criticidad compuesto) quedan soportados en su dimensión "
         "descriptiva — el sistema ya captura frecuencia, tiempo y coste, falta combinarlos en una métrica única.")

# ============================================================================
# SLIDE 19 — LIMITACIONES Y TRABAJO FUTURO
# ============================================================================
s = content("Limitaciones y trabajo futuro", "Prospectiva", "Conclusiones", 19)
card(s, ML, 2.15, (CW - 0.5) / 2, 4.25, fill=CARD, edge=AMBER)
txt(s, ML + 0.3, 2.4, (CW - 0.5) / 2 - 0.6, 0.4, [[R("Limitaciones", 15, INKTXT, True, F_HEAD)]])
bullets(s, ML + 0.3, 2.95, (CW - 0.5) / 2 - 0.6, 3.3, [
    "LLMs de 3B en CPU: latencia de inferencia considerable (cuello de botella del pipeline).",
    "Prototipo demostrador, no sistema de producción: sin autenticación ni control de acceso.",
    "Validado sobre datos sintéticos, no reales de empresa.",
    "OE6/OE7 cubiertos en su dimensión descriptiva, no como métrica compuesta.",
], size=12, gap=10, bcol=AMBER)
x2 = ML + (CW - 0.5) / 2 + 0.5
card(s, x2, 2.15, (CW - 0.5) / 2, 4.25, fill=CARD, edge=TEAL)
txt(s, x2 + 0.3, 2.4, (CW - 0.5) / 2 - 0.6, 0.4, [[R("Trabajo futuro", 15, INKTXT, True, F_HEAD)]])
bullets(s, x2 + 0.3, 2.95, (CW - 0.5) / 2 - 0.6, 3.3, [
    "Índice de criticidad compuesto (OE7) y serie temporal de incidencias (OE6).",
    "Webhook de subida (web, carpeta monitorizada o nube) para la ingesta.",
    "Notificación activa (Slack/email) ante inconsistencias de severidad alta.",
    "Más reglas de correlación del manual y alertas de Grafana sobre umbrales.",
], size=12, gap=10, bcol=TEAL)
notes(s, "PRESENTA: Víctor. ~50 s.\nSomos honestos con el alcance: el mayor límite es la latencia de los "
         "LLMs locales en CPU, y es un demostrador, no producción. La hoja de ruta es clara: cerrar el índice "
         "compuesto, añadir ingesta por webhook, notificaciones activas y más reglas de coherencia.")

# ============================================================================
# SLIDE 20 — CONCLUSIONES
# ============================================================================
s = content("Conclusiones", "Cierre", "Conclusiones", 20)
card(s, ML, 2.1, CW, 1.35, fill=INK)
rect(s, ML, 2.1, 0.12, 1.35, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
txt(s, ML + 0.35, 2.1, CW - 0.7, 1.35,
    [[R("OpsInsight Analytics demuestra que es posible pasar de documentación operativa heterogénea a "
        "inteligencia analítica accionable mediante un enfoque local, reproducible y trazable.",
        14.5, WHITE, True, F_HEAD)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
txt(s, ML, 3.75, CW, 0.35, [[R("Tres aportaciones diferenciales", 13, INKTXT, True, F_HEAD)]])
aportes = [
    ("Pipeline local end-to-end", "Arquitectura multiagente completa sin dependencias en la nube.", TEAL),
    ("Validación semántica", "Coherencia del dato contra el manual técnico, no solo estructural.", AMBER),
    ("Visualización trazable", "Del reporte a la decisión, con visibilidad completa del pipeline.", GREEN),
]
aw = (CW - 0.5 * 2) / 3
ax = ML
for tt, ds, col in aportes:
    card(s, ax, 4.2, aw, 1.7, fill=CARD, edge=col)
    txt(s, ax + 0.25, 4.4, aw - 0.5, 1.35,
        [[R(tt, 13.5, INKTXT, True, F_HEAD)], [R(ds, 11, MUTED, False, F_BODY)]],
        space_after=5, line_spacing=1.08)
    ax += aw + 0.5
card(s, ML, 6.1, CW, 0.5, fill=TEALSOFT, edge=TEAL)
txt(s, ML + 0.3, 6.1, CW - 0.6, 0.5,
    [[R("Los objetivos nucleares se alcanzaron y validaron de extremo a extremo; el resultado define una "
        "hoja de ruta clara hacia un producto profesional.", 11.5, TEAL_D, True, F_HEAD)]],
    anchor=MSO_ANCHOR.MIDDLE)
notes(s, "PRESENTA: Efrel. ~50 s.\nEn síntesis: materializamos una respuesta concreta al vacío del estado del "
         "arte. Tres aportaciones diferenciales — pipeline local completo, validación semántica y "
         "visualización trazable. Un MVP honesto que abre una hoja de ruta clara.")

# ============================================================================
# SLIDE 21 — CIERRE
# ============================================================================
s = slide()
rect(s, 0, 0, SW, SH, fill=INK)
rect(s, 0, 0, SW, 0.22, fill=TEAL)
rect(s, 0, SH - 0.22, SW, 0.22, fill=TEAL)
txt(s, ML, 2.35, CW, 1.2, [[R("Gracias por su atención", 40, WHITE, True, F_HEAD)]])
txt(s, ML, 3.5, CW, 0.7, [[R("¿Preguntas?", 22, TEAL, True, F_HEAD)]])
rect(s, ML + 0.03, 4.35, 2.4, 0.06, fill=TEAL)
txt(s, ML, 4.6, CW, 0.4,
    [[R("OpsInsight Analytics", 14, WHITE, True, F_HEAD),
      R("   ·   Máster en Análisis y Visualización de Datos Masivos · UNIR", 12, RGBColor(0xC6, 0xD3, 0xDD), False, F_BODY)]])
txt(s, ML, 5.15, CW, 0.4,
    [[R("Efrel Armando López Cáceres   ·   José Miguel Torres Cámara   ·   Víctor Romero Pardeiro",
        12.5, RGBColor(0xC6, 0xD3, 0xDD), False, F_BODY)]])
rect(s, ML, 5.9, 4.15, 0.5, fill=PANEL, line=TEAL_D, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.4)
txt(s, ML + 0.22, 5.9, 3.9, 0.5,
    [[R("RedProyectum", 11, TEAL, True, F_HEAD), R("  ·  Sopra Steria", 10.5, RGBColor(0xC6, 0xD3, 0xDD), False, F_BODY)]],
    anchor=MSO_ANCHOR.MIDDLE)
notes(s, "PRESENTA: los tres. Cierre y turno de preguntas del tribunal.\n"
         "Posibles preguntas: ¿por qué n8n y no Activepieces? ¿por qué LLM local? ¿cómo se evita la "
         "alucinación del agente? ¿escalabilidad a datos reales?")

prs.save(OUT)
print("OK ->", OUT)
print("Slides:", len(prs.slides._sldIdLst))
